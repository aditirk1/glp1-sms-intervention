"""Build the GoaLPost¹ walkthrough PowerPoint.

Run:  python docs/build_pptx.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent / "GoaLPost_Walkthrough.pptx"

INK = RGBColor(0x1A, 0x23, 0x32)
MUTED = RGBColor(0x5A, 0x65, 0x77)
PAPER = RGBColor(0xF7, 0xF4, 0xEF)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xD9, 0xD2, 0xC5)
TEAL = RGBColor(0x0F, 0x6B, 0x5C)
TEAL_SOFT = RGBColor(0xD8, 0xEF, 0xE9)
WARN = RGBColor(0xB4, 0x53, 0x09)
WARN_SOFT = RGBColor(0xFE, 0xF3, 0xC7)
ALERT = RGBColor(0x9F, 0x12, 0x39)
ALERT_SOFT = RGBColor(0xFF, 0xE4, 0xE6)
OK = RGBColor(0x16, 0x65, 0x34)
OK_SOFT = RGBColor(0xDC, 0xFC, 0xE7)
AMBER = RGBColor(0xC2, 0x41, 0x0C)
AMBER_SOFT = RGBColor(0xFF, 0xED, 0xD5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Avenir"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _set_run(run, *, size=18, bold=False, color=INK, font=FONT):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def write_block(shape, lines, *, default_size=18, default_color=INK):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True

    for index, item in enumerate(lines):
        if isinstance(item, str):
            text, kwargs = item, {}
        else:
            text, kwargs = item[0], (item[1] if len(item) > 1 else {})

        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = ""
        p.alignment = kwargs.get("align", PP_ALIGN.LEFT)
        p.space_after = Pt(kwargs.get("space_after", 6))
        p.level = kwargs.get("level", 0)
        run = p.add_run()
        run.text = text
        _set_run(
            run,
            size=kwargs.get("size", default_size),
            bold=kwargs.get("bold", False),
            color=kwargs.get("color", default_color),
            font=kwargs.get("font", FONT),
        )


def fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def rounded_rect(slide, left, top, width, height, fill_color=CARD, line=True):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line:
        shape.line.color.rgb = LINE
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    try:
        shape.adjustments[0] = 0.08
    except Exception:
        pass
    return shape


def accent_bar(slide):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.1), SLIDE_H
    )
    fill(bar, TEAL)
    return bar


def footer(slide, page, total):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(7.05), Inches(4), Inches(0.3))
    write_block(box, [("GoaLPost¹", {"size": 11, "bold": True, "color": INK})])
    box2 = slide.shapes.add_textbox(Inches(11.2), Inches(7.05), Inches(1.6), Inches(0.3))
    write_block(
        box2,
        [(f"{page} / {total}", {"size": 11, "color": MUTED, "align": PP_ALIGN.RIGHT})],
    )
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(6.95), Inches(12.25), Pt(1)
    )
    fill(line, LINE)


def eyebrow(slide, text, left=Inches(0.55), top=Inches(0.38)):
    box = slide.shapes.add_textbox(left, top, Inches(11), Inches(0.3))
    write_block(
        box,
        [(text.upper(), {"size": 11, "bold": True, "color": TEAL, "space_after": 0})],
    )


def title(slide, text, left=Inches(0.55), top=Inches(0.72), width=Inches(12.2), size=28):
    box = slide.shapes.add_textbox(left, top, width, Inches(1.0))
    write_block(
        box,
        [(text, {"size": size, "bold": True, "color": INK, "space_after": 0})],
    )


def body(slide, text, left=Inches(0.55), top=Inches(1.7), width=Inches(12), size=16, color=MUTED):
    box = slide.shapes.add_textbox(left, top, width, Inches(0.85))
    write_block(box, [(text, {"size": size, "color": color, "space_after": 0})])


def card(slide, left, top, width, height, *, fill_color=CARD):
    return rounded_rect(slide, left, top, width, height, fill_color=fill_color)


def pill(slide, left, top, text, bg, fg):
    w = Inches(0.1) * max(len(text), 4) + Inches(0.5)
    shape = rounded_rect(slide, left, top, w, Inches(0.3), fill_color=bg, line=False)
    try:
        shape.adjustments[0] = 0.5
    except Exception:
        pass
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(2)
    p.text = ""
    run = p.add_run()
    run.text = text.upper()
    _set_run(run, size=9, bold=True, color=fg)
    return shape


def callout(slide, left, top, width, height, text, *, warn=False):
    bg = WARN_SOFT if warn else TEAL_SOFT
    accent = WARN if warn else TEAL
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.07), height)
    fill(bar, accent)
    panel = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left + Inches(0.07), top, width - Inches(0.07), height
    )
    fill(panel, bg)
    box = slide.shapes.add_textbox(
        left + Inches(0.28), top + Inches(0.14), width - Inches(0.45), height - Inches(0.2)
    )
    write_block(box, [(text, {"size": 14, "color": INK, "space_after": 0})])


def blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    fill(bg, PAPER)
    spTree = slide.shapes._spTree
    sp = bg._element
    spTree.remove(sp)
    spTree.insert(2, sp)
    accent_bar(slide)
    return slide


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    total = 17
    page = 0

    def next_page():
        nonlocal page
        page += 1
        return page

    # 1 Title
    s = blank_slide(prs)
    p = next_page()
    o1 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.8), Inches(-1.0), Inches(4.5), Inches(4.5))
    fill(o1, TEAL_SOFT)
    o2 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.2), Inches(4.8), Inches(4), Inches(4))
    fill(o2, RGBColor(0xE4, 0xEC, 0xF5))
    eyebrow(s, "Project walkthrough", top=Inches(2.15))
    title(s, "GoaLPost¹", top=Inches(2.55), size=48)
    body(
        s,
        "An automated SMS check-in system for patients on GLP-1 therapy.\n"
        "It identifies who is at risk of stopping treatment and routes the right follow-up.",
        top=Inches(3.7),
        width=Inches(10.5),
        size=18,
    )
    body(
        s,
        "Prepared for a non-technical audience. No coding background assumed.",
        top=Inches(5.0),
        size=14,
    )
    footer(s, p, total)

    # 2 Problem
    s = blank_slide(prs)
    p = next_page()
    eyebrow(s, "Context")
    title(s, "Retention on GLP-1 therapy is a known problem")
    card(s, Inches(0.55), Inches(2.05), Inches(5.9), Inches(2.55))
    pill(s, Inches(0.8), Inches(2.3), "Clinical pattern", WARN_SOFT, WARN)
    box = s.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(5.4), Inches(1.55))
    write_block(
        box,
        [
            ("Many patients discontinue within the first months", {"size": 16, "bold": True, "space_after": 8}),
            (
                "Common drivers include gastrointestinal side effects, stalled weight loss, "
                "coverage or cost barriers, and loss of contact with the care team.",
                {"size": 14, "color": MUTED},
            ),
        ],
    )
    card(s, Inches(6.8), Inches(2.05), Inches(5.9), Inches(2.55))
    pill(s, Inches(7.05), Inches(2.3), "Operational gap", TEAL_SOFT, TEAL)
    box = s.shapes.add_textbox(Inches(7.05), Inches(2.8), Inches(5.4), Inches(1.55))
    write_block(
        box,
        [
            ("Manual outreach does not scale across a full panel", {"size": 16, "bold": True, "space_after": 8}),
            (
                "Care teams cannot personally text every patient each week. "
                "By the time someone reports a problem, they may already have stopped filling.",
                {"size": 14, "color": MUTED},
            ),
        ],
    )
    callout(
        s,
        Inches(0.55),
        Inches(5.0),
        Inches(12.2),
        Inches(1.3),
        "GoaLPost¹ monitors a patient panel on a schedule, estimates who is drifting, "
        "attributes a likely reason, and either sends a targeted message or creates "
        "a work item for a clinician.",
    )
    footer(s, p, total)

    # 3 Patient view
    s = blank_slide(prs)
    p = next_page()
    eyebrow(s, "Patient experience")
    title(s, "Check-ins are a single SMS with three reply options")
    card(s, Inches(0.55), Inches(2.0), Inches(8.8), Inches(3.5))
    box = s.shapes.add_textbox(Inches(0.9), Inches(2.3), Inches(8.1), Inches(3.0))
    write_block(
        box,
        [
            (
                "Hi Maria, week 12 check-in from GoaLPost¹.\nHow is your GLP-1 journey going?",
                {"size": 18, "bold": True, "space_after": 16},
            ),
            ("1  Going well", {"size": 16, "space_after": 8}),
            ("2  Having side effects", {"size": 16, "space_after": 8}),
            ("3  Not seeing results", {"size": 16, "space_after": 8}),
        ],
    )
    body(
        s,
        "Patients do not need an app. The clinical decision support sits with the care team.",
        top=Inches(5.85),
        size=15,
    )
    footer(s, p, total)

    # 4 Big idea
    s = blank_slide(prs)
    p = next_page()
    eyebrow(s, "Design")
    title(s, "Scheduling and follow-up run without a manual send button")
    body(
        s,
        "A single decision process runs on a timer. On each pass it answers three questions:",
        top=Inches(1.7),
        size=16,
    )
    cards = [
        ("1", "Who is due for contact?", "Contact frequency depends on risk. Higher-risk patients are checked more often."),
        ("2", "How high is the dropout risk?", "A score estimates the chance the patient stops therapy within roughly 90 days."),
        ("3", "What action is appropriate?", "Send a targeted message, escalate to a clinician, or wait until the next scheduled check."),
    ]
    for i, (n, h, t) in enumerate(cards):
        left = Inches(0.55) + Inches(i * 4.15)
        card(s, left, Inches(2.55), Inches(3.9), Inches(2.7))
        pill(s, left + Inches(0.25), Inches(2.8), n, TEAL_SOFT, TEAL)
        box = s.shapes.add_textbox(left + Inches(0.25), Inches(3.3), Inches(3.4), Inches(1.7))
        write_block(
            box,
            [
                (h, {"size": 15, "bold": True, "space_after": 8}),
                (t, {"size": 13, "color": MUTED}),
            ],
        )
    callout(
        s,
        Inches(0.55),
        Inches(5.55),
        Inches(12.2),
        Inches(0.9),
        "The live product and the multi-month simulation call the same decision process, "
        "so the demo exercises production logic rather than a separate prototype path.",
    )
    footer(s, p, total)

    # 5 Pipeline
    s = blank_slide(prs)
    p = next_page()
    eyebrow(s, "Pipeline")
    title(s, "What happens on each scheduled run")
    steps = [
        ("1", "Timer fires", "Hourly in production, or via an external cron job on hosts that sleep when idle"),
        ("2", "Select due patients", "Anyone whose next check-in time has passed"),
        ("3", "Score risk", "Assign green, amber, or red, plus a likely barrier"),
        ("4", "Apply policy", "Send a message and/or create a care-team task"),
        ("5", "Reschedule", "Set the next check-in based on the updated risk tier"),
    ]
    for i, (n, h, t) in enumerate(steps):
        left = Inches(0.4) + Inches(i * 2.55)
        card(s, left, Inches(2.15), Inches(2.35), Inches(3.15))
        box = s.shapes.add_textbox(left + Inches(0.15), Inches(2.35), Inches(2.05), Inches(2.75))
        write_block(
            box,
            [
                (n, {"size": 11, "bold": True, "color": TEAL, "space_after": 6}),
                (h, {"size": 14, "bold": True, "space_after": 8}),
                (t, {"size": 12, "color": MUTED}),
            ],
        )
        if i < 4:
            arrow = s.shapes.add_textbox(
                left + Inches(2.2), Inches(3.4), Inches(0.4), Inches(0.4)
            )
            write_block(arrow, [(">", {"size": 18, "color": MUTED, "align": PP_ALIGN.CENTER})])
    body(
        s,
        "When a patient replies, the system scores that reply immediately, applies the same policy, "
        "and updates their next check-in date.",
        top=Inches(5.6),
        size=14,
    )
    footer(s, p, total)

    # 6 Cadence
    s = blank_slide(prs)
    p = next_page()
    eyebrow(s, "Contact frequency")
    title(s, "Higher risk shortens the time between check-ins")
    tiers = [
        ("Green", "Every 14 days", "Patient appears stable; outreach stays light.", OK_SOFT, OK),
        ("Amber", "Every 7 days", "Elevated concern; maintain weekly contact.", AMBER_SOFT, AMBER),
        ("Red", "Every 3 days", "High dropout risk; keep a short follow-up loop.", ALERT_SOFT, ALERT),
    ]
    for i, (name, freq, desc, bg, fg) in enumerate(tiers):
        left = Inches(0.55) + Inches(i * 4.15)
        card(s, left, Inches(2.1), Inches(3.9), Inches(2.45))
        pill(s, left + Inches(0.25), Inches(2.35), name, bg, fg)
        box = s.shapes.add_textbox(left + Inches(0.25), Inches(2.9), Inches(3.4), Inches(1.4))
        write_block(
            box,
            [
                (freq, {"size": 20, "bold": True, "space_after": 8}),
                (desc, {"size": 14, "color": MUTED}),
            ],
        )
    callout(
        s,
        Inches(0.55),
        Inches(5.0),
        Inches(12.2),
        Inches(1.3),
        "During the first four weeks of therapy, the gap between check-ins is capped at seven days "
        "for every patient. Early dose titration is when gastrointestinal side effects and early "
        "discontinuation are most concentrated.",
        warn=True,
    )
    footer(s, p, total)

    # 7 Risk score
    s = blank_slide(prs)
    p = next_page()
    eyebrow(s, "Risk scoring")
    title(s, "The risk score is a forecast of discontinuation")
    body(
        s,
        "It is a number between 0 and 1 describing how likely the patient is to stop therapy "
        "within about 90 days. You can think of it like a weather forecast for retention.",
        top=Inches(1.7),
        size=15,
    )
    card(s, Inches(0.55), Inches(2.55), Inches(6.0), Inches(3.7))
    box = s.shapes.add_textbox(Inches(0.85), Inches(2.8), Inches(5.4), Inches(3.2))
    write_block(
        box,
        [
            ("Inputs to the forecast", {"size": 16, "bold": True, "space_after": 10}),
            ("Time on therapy", {"size": 14, "color": MUTED, "space_after": 4}),
            ("Whether weight appears stuck (inferred from replies in the demo)", {"size": 14, "color": MUTED, "space_after": 4}),
            ("Reported side effects", {"size": 14, "color": MUTED, "space_after": 4}),
            ("Insurance type and income band", {"size": 14, "color": MUTED, "space_after": 4}),
            ("Repeated replies of \"not seeing results\"", {"size": 14, "color": MUTED, "space_after": 4}),
            ("Number of check-ins left unanswered", {"size": 14, "bold": True, "color": INK, "space_after": 4}),
        ],
    )
    card(s, Inches(6.85), Inches(2.55), Inches(5.9), Inches(3.7))
    box = s.shapes.add_textbox(Inches(7.15), Inches(2.8), Inches(5.3), Inches(3.2))
    write_block(
        box,
        [
            ("Outputs shown to the care team", {"size": 16, "bold": True, "space_after": 10}),
            ("A probability converted into green, amber, or red", {"size": 14, "color": MUTED, "space_after": 10}),
            ("A likely barrier: the main reason the forecast worsened", {"size": 14, "color": MUTED, "space_after": 6}),
            ("Plateau, side effects, or cost", {"size": 15, "bold": True, "space_after": 12}),
            (
                "An explanation layer highlights which factors mattered most for that patient. "
                "The dashboard presents the conclusion; the underlying statistics stay behind the scenes.",
                {"size": 13, "color": MUTED},
            ),
        ],
    )
    footer(s, p, total)

    # 8 Playbook
    s = blank_slide(prs)
    p = next_page()
    eyebrow(s, "Intervention policy")
    title(s, "Follow-up rules are written as an explicit playbook")
    body(
        s,
        "The policy is stored as a readable list so a clinician can review what the system will do.",
        top=Inches(1.65),
        size=14,
    )
    rows = [
        ("Plateau", "Patient reports \"not seeing results\" with elevated risk: send a message grounded in published trial evidence about non-scale benefits."),
        ("Side effects", "Reply 2: send practical gastrointestinal guidance. Two reports within a month: escalate to the care team."),
        ("Cost", "Coverage or cost flagged as the driver: send a navigation message and open a benefits-check task."),
        ("Silence", "One missed check-in: send a short reminder. Two missed check-ins: open an outreach call task."),
        ("Sustained red", "Red on two consecutive scorings: open a nurse call task."),
        ("Doing well", "Reply 1: send a brief acknowledgment and allow a longer interval before the next check-in."),
    ]
    for i, (label, desc) in enumerate(rows):
        top = Inches(2.1) + Inches(i * 0.7)
        card(s, Inches(0.55), top, Inches(12.2), Inches(0.6))
        label_box = s.shapes.add_textbox(Inches(0.75), top + Inches(0.12), Inches(2.1), Inches(0.4))
        write_block(label_box, [(label, {"size": 13, "bold": True, "color": TEAL, "space_after": 0})])
        desc_box = s.shapes.add_textbox(Inches(2.95), top + Inches(0.1), Inches(9.5), Inches(0.45))
        write_block(desc_box, [(desc, {"size": 12, "color": MUTED, "space_after": 0})])
    footer(s, p, total)

    # 9 Guardrails
    s = blank_slide(prs)
    p = next_page()
    eyebrow(s, "Safety constraints")
    title(s, "Limits prevent message spam and queue overload")
    card(s, Inches(0.55), Inches(2.0), Inches(6.0), Inches(3.2))
    box = s.shapes.add_textbox(Inches(0.85), Inches(2.25), Inches(5.4), Inches(2.7))
    write_block(
        box,
        [
            ("Outbound messages", {"size": 16, "bold": True, "space_after": 10}),
            ("At most two messages per patient in any rolling seven-day window", {"size": 13, "color": MUTED, "space_after": 6}),
            ("No repeated message body within fourteen days", {"size": 13, "color": MUTED, "space_after": 6}),
            ("Sends only between 9 a.m. and 8 p.m.", {"size": 13, "color": MUTED, "space_after": 6}),
            ("At most one message per decision pass", {"size": 13, "color": MUTED, "space_after": 6}),
        ],
    )
    card(s, Inches(6.85), Inches(2.0), Inches(5.9), Inches(3.2))
    box = s.shapes.add_textbox(Inches(7.15), Inches(2.25), Inches(5.3), Inches(2.7))
    write_block(
        box,
        [
            ("Care-team tasks", {"size": 16, "bold": True, "space_after": 10}),
            ("A second task of the same type is not opened while one remains open", {"size": 13, "color": MUTED, "space_after": 6}),
            ("A cooldown applies after a task is closed so standing conditions do not refill the queue", {"size": 13, "color": MUTED, "space_after": 6}),
            ("Open work is ranked by priority, then by current risk", {"size": 13, "color": MUTED, "space_after": 6}),
        ],
    )
    callout(
        s,
        Inches(0.55),
        Inches(5.5),
        Inches(12.2),
        Inches(0.9),
        "In a stress test of 1,000 patients over 26 weeks, the weekly message cap was never "
        "exceeded and no messages were sent outside the allowed hours.",
    )
    footer(s, p, total)

    # 10 Silence
    s = blank_slide(prs)
    p = next_page()
    eyebrow(s, "Non-response")
    title(s, "Unanswered check-ins are treated as a clinical signal")
    card(s, Inches(0.55), Inches(2.1), Inches(5.9), Inches(3.0))
    pill(s, Inches(0.8), Inches(2.35), "Earlier approach", WARN_SOFT, WARN)
    box = s.shapes.add_textbox(Inches(0.8), Inches(2.9), Inches(5.4), Inches(1.9))
    write_block(
        box,
        [
            ("Quiet patients were easy to overlook", {"size": 16, "bold": True, "space_after": 8}),
            (
                "Therapy week only advanced when someone replied. A patient who never "
                "answered could remain recorded at week zero and look artificially low risk.",
                {"size": 14, "color": MUTED},
            ),
        ],
    )
    card(s, Inches(6.8), Inches(2.1), Inches(5.9), Inches(3.0))
    pill(s, Inches(7.05), Inches(2.35), "Current approach", TEAL_SOFT, TEAL)
    box = s.shapes.add_textbox(Inches(7.05), Inches(2.9), Inches(5.4), Inches(1.9))
    write_block(
        box,
        [
            ("Missed check-ins raise risk and create work", {"size": 16, "bold": True, "space_after": 8}),
            (
                "Time on therapy advances from the calendar. Unanswered prompts feed the "
                "risk score. Prolonged silence opens outreach work and can mark the patient "
                "as lost to follow-up.",
                {"size": 14, "color": MUTED},
            ),
        ],
    )
    body(
        s,
        "This matches clinical intuition: patients who stop answering often stop refilling.",
        top=Inches(5.5),
        size=15,
    )
    footer(s, p, total)

    # 11 Dashboard
    s = blank_slide(prs)
    p = next_page()
    eyebrow(s, "Care team interface")
    title(s, "The dashboard opens on work that needs a person")
    items = [
        ("Work queue", "Open tasks ordered by priority and risk. Staff can acknowledge or resolve each item."),
        ("Cohort summary", "Distribution across risk tiers, response rate, and retention over time."),
        ("Patient roster", "Filter by risk tier, silence, or status, and see who is next due for contact."),
    ]
    for i, (h, t) in enumerate(items):
        left = Inches(0.55) + Inches(i * 4.15)
        card(s, left, Inches(2.15), Inches(3.9), Inches(2.5))
        box = s.shapes.add_textbox(left + Inches(0.25), Inches(2.45), Inches(3.4), Inches(1.9))
        write_block(
            box,
            [
                (h, {"size": 16, "bold": True, "space_after": 10}),
                (t, {"size": 14, "color": MUTED}),
            ],
        )
    callout(
        s,
        Inches(0.55),
        Inches(5.1),
        Inches(12.2),
        Inches(1.2),
        "Clinician time is reserved for cases where a text is unlikely to be enough. "
        "The product surface is a prioritized work queue, not a full transcript of every message.",
    )
    footer(s, p, total)

    # 12 Simulation
    s = blank_slide(prs)
    p = next_page()
    eyebrow(s, "Evaluation")
    title(s, "We compared two arms over 26 simulated weeks")
    card(s, Inches(0.55), Inches(2.0), Inches(5.9), Inches(2.4))
    pill(s, Inches(0.8), Inches(2.25), "Control", AMBER_SOFT, AMBER)
    box = s.shapes.add_textbox(Inches(0.8), Inches(2.75), Inches(5.4), Inches(1.4))
    write_block(
        box,
        [
            ("Check-ins and scoring only", {"size": 17, "bold": True, "space_after": 8}),
            (
                "Patients still receive scheduled prompts and risk scores. "
                "Targeted interventions and escalations are turned off.",
                {"size": 13, "color": MUTED},
            ),
        ],
    )
    card(s, Inches(6.8), Inches(2.0), Inches(5.9), Inches(2.4))
    pill(s, Inches(7.05), Inches(2.25), "Intervention", TEAL_SOFT, TEAL)
    box = s.shapes.add_textbox(Inches(7.05), Inches(2.75), Inches(5.4), Inches(1.4))
    write_block(
        box,
        [
            ("Full playbook enabled", {"size": 17, "bold": True, "space_after": 8}),
            (
                "Same check-ins, plus plateau, side-effect, and cost messages, "
                "and care-team tasks where the policy requires them.",
                {"size": 13, "color": MUTED},
            ),
        ],
    )
    body(
        s,
        "Each arm used 1,000 simulated patients for 26 weeks, advancing one day at a time "
        "through the same decision process used in production.",
        top=Inches(4.65),
        size=14,
        color=INK,
    )
    callout(
        s,
        Inches(0.55),
        Inches(5.35),
        Inches(12.2),
        Inches(1.1),
        "Important limitation: how much each message type helps is an assumption declared "
        "in the simulation code. The run shows that scheduling, guardrails, and escalation "
        "behave correctly at cohort scale. It does not establish clinical effect sizes. "
        "Those require a prospective pilot.",
        warn=True,
    )
    footer(s, p, total)

    # 13 Results
    s = blank_slide(prs)
    p = next_page()
    eyebrow(s, "Results")
    title(s, "Outcomes under the declared simulation assumptions")
    card(s, Inches(0.55), Inches(2.0), Inches(5.9), Inches(2.55))
    pill(s, Inches(0.8), Inches(2.25), "Control", AMBER_SOFT, AMBER)
    box = s.shapes.add_textbox(Inches(0.8), Inches(2.7), Inches(5.4), Inches(1.6))
    write_block(
        box,
        [
            ("40%", {"size": 44, "bold": True, "color": ALERT, "space_after": 6}),
            ("still marked active at week 26", {"size": 14, "color": MUTED}),
        ],
    )
    card(s, Inches(6.8), Inches(2.0), Inches(5.9), Inches(2.55))
    pill(s, Inches(7.05), Inches(2.25), "Intervention", TEAL_SOFT, TEAL)
    box = s.shapes.add_textbox(Inches(7.05), Inches(2.7), Inches(5.4), Inches(1.6))
    write_block(
        box,
        [
            ("62%", {"size": 44, "bold": True, "color": TEAL, "space_after": 6}),
            ("still marked active, a difference of about 22 points", {"size": 14, "color": MUTED}),
        ],
    )
    metrics = [
        ("Weekly message cap", "Maximum of two; zero violations"),
        ("Send window", "No messages outside 9 a.m. to 8 p.m."),
        ("Silent patients", "None remained frozen at week zero"),
    ]
    for i, (h, t) in enumerate(metrics):
        left = Inches(0.55) + Inches(i * 4.15)
        card(s, left, Inches(4.9), Inches(3.9), Inches(1.4))
        box = s.shapes.add_textbox(left + Inches(0.25), Inches(5.1), Inches(3.4), Inches(1.0))
        write_block(
            box,
            [
                (h, {"size": 13, "bold": True, "space_after": 4}),
                (t, {"size": 12, "color": MUTED}),
            ],
        )
    footer(s, p, total)

    # 14 Demo path
    s = blank_slide(prs)
    p = next_page()
    eyebrow(s, "Live demonstration")
    title(s, "Suggested order when presenting the product")
    steps = [
        "Open the dashboard and start with the work queue rather than the full patient table.",
        "Acknowledge and resolve one task to show how staff close work items.",
        "Run the scheduler once to contact due patients. Running it again should contact nobody new.",
        "Simulate a patient reply of 3 and review the risk score, barrier, and plateau message.",
        "Optionally simulate a reply of 2 to show the side-effect path.",
        "Show the control versus intervention chart, and state the simulation limitation clearly.",
    ]
    for i, text in enumerate(steps):
        top = Inches(1.9) + Inches(i * 0.72)
        circ = s.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.55), top + Inches(0.05), Inches(0.4), Inches(0.4)
        )
        fill(circ, TEAL)
        nbox = s.shapes.add_textbox(Inches(0.55), top + Inches(0.1), Inches(0.4), Inches(0.32))
        write_block(
            nbox,
            [(str(i + 1), {"size": 13, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER, "space_after": 0})],
        )
        card(s, Inches(1.15), top, Inches(11.6), Inches(0.58))
        box = s.shapes.add_textbox(Inches(1.4), top + Inches(0.12), Inches(11.1), Inches(0.38))
        write_block(box, [(text, {"size": 13, "color": INK, "space_after": 0})])
    footer(s, p, total)

    # 15 Summary
    s = blank_slide(prs)
    p = next_page()
    eyebrow(s, "Summary")
    title(s, "What this project delivers")
    items = [
        "An automated check-in schedule that tightens as dropout risk rises",
        "A risk forecast with an attributed reason: plateau, side effects, cost, or silence",
        "An explicit playbook for messages and clinician escalation",
        "Operational limits on message volume and task creation",
        "A care-team work queue as the primary clinical interface",
        "A simulation harness that runs the same decision process over months of virtual time",
    ]
    for i, text in enumerate(items):
        top = Inches(1.95) + Inches(i * 0.7)
        dot = s.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.7), top + Inches(0.16), Inches(0.16), Inches(0.16)
        )
        fill(dot, TEAL)
        box = s.shapes.add_textbox(Inches(1.15), top, Inches(11.5), Inches(0.55))
        write_block(box, [(text, {"size": 16, "color": INK, "space_after": 0})])
    footer(s, p, total)

    # 16 Round 2
    s = blank_slide(prs)
    p = next_page()
    eyebrow(s, "Next phase")
    title(s, "Priorities for Round 2")
    card(s, Inches(0.55), Inches(2.0), Inches(6.0), Inches(3.5))
    pill(s, Inches(0.8), Inches(2.25), "Near term", TEAL_SOFT, TEAL)
    box = s.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(5.4), Inches(2.5))
    write_block(
        box,
        [
            ("Connect to real messaging and clinical review", {"size": 16, "bold": True, "space_after": 10}),
            ("Send and receive live SMS through Twilio on a test number", {"size": 13, "color": MUTED, "space_after": 6}),
            ("Review green, amber, and red thresholds with a clinician", {"size": 13, "color": MUTED, "space_after": 6}),
            ("Replace assumed message effect sizes with literature or advisor estimates", {"size": 13, "color": MUTED, "space_after": 6}),
            ("Use measured weight or lab data when available instead of reply proxies", {"size": 13, "color": MUTED, "space_after": 6}),
        ],
    )
    card(s, Inches(6.85), Inches(2.0), Inches(5.9), Inches(3.5))
    pill(s, Inches(7.15), Inches(2.25), "Then", WARN_SOFT, WARN)
    box = s.shapes.add_textbox(Inches(7.15), Inches(2.8), Inches(5.3), Inches(2.5))
    write_block(
        box,
        [
            ("Prepare for a small pilot", {"size": 16, "bold": True, "space_after": 10}),
            ("Add dashboard authentication and an audit trail for resolved tasks", {"size": 13, "color": MUTED, "space_after": 6}),
            ("Define a clinic endpoint such as 90-day persistence before launch", {"size": 13, "color": MUTED, "space_after": 6}),
            ("Keep improving operations rather than chasing model complexity on synthetic labels", {"size": 13, "color": MUTED, "space_after": 6}),
        ],
    )
    callout(
        s,
        Inches(0.55),
        Inches(5.8),
        Inches(12.2),
        Inches(0.7),
        "Round 1 established the system design. Round 2 should establish contact with "
        "real messaging, clinical cadence, and a pilot plan.",
    )
    footer(s, p, total)

    # 17 Close
    s = blank_slide(prs)
    p = next_page()
    o1 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.8), Inches(-1.0), Inches(4.5), Inches(4.5))
    fill(o1, TEAL_SOFT)
    o2 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.2), Inches(4.8), Inches(4), Inches(4))
    fill(o2, RGBColor(0xE4, 0xEC, 0xF5))
    eyebrow(s, "Discussion", top=Inches(2.2))
    title(s, "Questions", top=Inches(2.6), size=42)
    body(
        s,
        "We can walk through the live dashboard, the simulation results, or the Round 2 plan.",
        top=Inches(3.75),
        size=17,
    )
    extras = [
        ("Dashboard", "localhost:8501"),
        ("API", "localhost:8000"),
        ("This file", "docs/GoaLPost_Walkthrough.pptx"),
    ]
    for i, (h, t) in enumerate(extras):
        left = Inches(0.55) + Inches(i * 4.15)
        card(s, left, Inches(4.75), Inches(3.9), Inches(1.35))
        box = s.shapes.add_textbox(left + Inches(0.25), Inches(4.95), Inches(3.4), Inches(0.95))
        write_block(
            box,
            [
                (h, {"size": 14, "bold": True, "space_after": 4}),
                (t, {"size": 12, "color": MUTED}),
            ],
        )
    footer(s, p, total)

    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    build()
